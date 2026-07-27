import base64
import requests
from pptx import Presentation
from pptx.util import Inches, Pt

def download_mermaid(text, filename):
    encoded = base64.urlsafe_b64encode(text.encode('utf-8')).decode('utf-8')
    url = f"https://mermaid.ink/img/{encoded}"
    res = requests.get(url, timeout=10)
    if res.status_code == 200:
        with open(filename, 'wb') as f:
            f.write(res.content)
        print(f"Saved {filename}")
    else:
        print(f"Failed to save {filename}, status code: {res.status_code}")

print("Downloading Mermaid UML Diagrams as images...")

arch = '''graph TD;
    User-->|Access Web UI|Frontend[Frontend HTML/CSS/JS];
    Frontend-->|Login / View History|SpringBoot[Spring Boot Backend];
    Frontend-->|Send Text / Request URL analysis|SpringBoot;
    SpringBoot-->|Read & Write Encrypted Data|MySQL[(MySQL DB)];
    SpringBoot-->|Forward Payload Data|Flask[Python Flask API];
    Flask-->|Download HTML Text|BS4[BeautifulSoup4 Web Scraper];
    Flask-->|Query Search Terms|GoogleNews[Google News RSS];
    Flask-->|Predict Probability|MLModel(Scikit-Learn ML Model);
    MLModel-->|Return JSON Metrics|SpringBoot;
    SpringBoot-->|Return JSON Dashboard Payload|Frontend;'''

usecase = '''graph TD;
    User((User))
    User --> UC1("Register via OTP")
    User --> UC2("Input News Text/URL")
    User --> UC3("View Truth Prediction")
    User --> UC4("Review Line-by-Line")
    User --> UC5("Manage History")'''

sequence = '''sequenceDiagram
    actor User
    participant Browser
    participant SpringBoot
    participant FlaskAPI
    participant GoogleNews
    participant DB

    User->>Browser: Paste News URL and Detect
    Browser->>SpringBoot: POST /api/news/detect
    SpringBoot->>FlaskAPI: POST /predict (text=URL)
    FlaskAPI->>FlaskAPI: Scrape URL content
    FlaskAPI->>GoogleNews: Search for matching keywords
    GoogleNews-->>FlaskAPI: Return related count
    FlaskAPI->>FlaskAPI: ML Model predicts Fake/Real
    FlaskAPI-->>SpringBoot: Return JSON result
    SpringBoot->>DB: Encrypt data & Save History
    SpringBoot-->>Browser: Return HTTP 200
    Browser-->>User: Render Result UI'''

download_mermaid(arch, 'arch.png')
download_mermaid(usecase, 'usecase.png')
download_mermaid(sequence, 'sequence.png')

print("Generating PowerPoint...")
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AI Fake News Detector"
subtitle.text = "A Full-Stack Application for Misinformation Detection\n\nBy Your Name/Team"

def add_bullet_slide(prs, title_text, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = bullets[0]
    
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        if bullet.startswith("  -") or bullet.startswith("    -"):
            p.level = 1

def add_image_slide(prs, title_text, img_path):
    slide_layout = prs.slide_layouts[5] # blank slide with title
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    try:
        # Add image, center it roughly
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))
    except Exception as e:
        print(f"Could not add image {img_path}: {e}")

# Add all the content slides
add_bullet_slide(prs, "Abstract", [
    "The rapid spread of misinformation online demands an automated, reliable detection mechanism.",
    "This project introduces a full-stack AI Fake News Detector utilizing a Natural Language Processing (NLP) pipeline.",
    "It combines Machine Learning prediction with real-time heuristic Google News verification to provide a granular, sentence-level breakdown analysis of articles and URLs."
])

add_bullet_slide(prs, "Introduction", [
    "Problem Statement: Fake news negatively impacts society, politics, and public opinion.",
    "Solution: A user-friendly web dashboard where users can paste raw text or direct news URLs to instantly assess credibility.",
    "Key Features:",
    "  - NLP-powered truth prediction.",
    "  - Sentence-by-sentence analysis isolating false vs. verified claims.",
    "  - Live integration with Google News scraping.",
    "  - Secure user authentication and History tracking."
])

add_bullet_slide(prs, "System Requirements", [
    "Hardware Requirements:",
    "  - Processor: Intel Core i5 or equivalent (Minimum)",
    "  - RAM: 8 GB (16 GB Recommended)",
    "  - Storage: 10 GB free space",
    "  - Network: Active Internet Connection (for live scraping)"
])

add_bullet_slide(prs, "Used Softwares & Required Frameworks", [
    "Frontend Technologies: HTML5, CSS3, Vanilla JavaScript",
    "Backend Technologies: Java Spring Boot, Python Flask",
    "Database: MySQL Server (AES Encryption)",
    "Machine Learning: Scikit-Learn, Pandas",
    "Data Scraping Tools: BeautifulSoup4, Requests"
])

add_bullet_slide(prs, "Methodology & Background", [
    "Data Source Training: Model trained offline on curated datasets of verified real vs. fake news.",
    "Core NLP Pipeline: TF-IDF Vectorizer with Logistic Regression.",
    "Hybrid Approach:",
    "  - Automater Web Scraper: Python BeautifulSoup4 isolates pure text from raw news URLs.",
    "  - Live Fact Verification: Pings Google News RSS to cross-reference search items."
])

# Add the 3 UML Image slides
add_image_slide(prs, "UML: System Architecture", "arch.png")
add_image_slide(prs, "UML: Use Case Diagram", "usecase.png")
add_image_slide(prs, "UML: Sequence Diagram", "sequence.png")

# Final screen
add_bullet_slide(prs, "Final Outputs (Screenshots required)", [
    "Please drop screenshots of your running application on the next slides for:",
    "  - Login Screen (Double-Password & OTP features)",
    "  - Main Dashboard",
    "  - Fake / Real Result Screen",
    "  - Sentence Breakdown Analysis highlighting verified claims",
    "  - User History Page showing previous detected URLs"
])

add_bullet_slide(prs, "Conclusion", [
    "This application intelligently bridges Machine Learning with live web-scraping to offer high-accuracy fake news detection.",
    "It provides transparent feedback on exactly which sentences are problematic, offering fact-checking context beyond a simple 'Fake/Real' label."
])

add_bullet_slide(prs, "Questions?", [
    "Thank you for listening!"
])

output_file = '../AI_Fake_News_Detector_Presentation.pptx'
prs.save(output_file)
print(f"Successfully generated full PPT at {output_file}")
